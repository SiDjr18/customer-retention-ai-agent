"""
ReportService — McKinsey-grade PDF, Excel, and PowerPoint reports.

Dependencies (all free / open-source):
  PDF   → reportlab
  Excel → openpyxl
  PPT   → python-pptx
"""

from __future__ import annotations

import os
from datetime import datetime, timezone
from typing import List

from app.config import settings
from app.schemas.report_schema import ReportRequest, ReportResponse
from app.services.dataset_service import DatasetService
from app.services.recommendation_service import RecommendationService

_TS_FMT = "%Y%m%d_%H%M%S"

NAVY   = "#003366"
TEAL   = "#0066CC"
GOLD   = "#F5A623"
LIGHT  = "#EEF4FB"
WHITE  = "#FFFFFF"
GREY   = "#F7F9FC"
RED    = "#C0392B"
GREEN  = "#27AE60"
TEXT   = "#1A1A2E"
MUTED  = "#6B7280"


def _ts() -> str:
    return datetime.now(timezone.utc).strftime(_TS_FMT)

def _ensure_reports_dir() -> str:
    os.makedirs(settings.REPORTS_DIR, exist_ok=True)
    return settings.REPORTS_DIR

def _risk_label(pct: float) -> str:
    if pct >= 25: return "CRITICAL"
    if pct >= 20: return "HIGH"
    if pct >= 15: return "ELEVATED"
    return "MODERATE"

def _fmt_inr(v: float) -> str:
    if v >= 10_000_000: return f"₹{v/10_000_000:.1f} Cr"
    if v >= 100_000:    return f"₹{v/100_000:.1f} L"
    return f"₹{v:,.0f}"


class ReportService:
    def __init__(self) -> None:
        self._ds = DatasetService()
        self._rs = RecommendationService()

    # ──────────────────────────────────────────────────────────────────────────
    # PDF — McKinsey-style executive report
    # ──────────────────────────────────────────────────────────────────────────

    def generate_pdf(self, req: ReportRequest) -> ReportResponse:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            HRFlowable, PageBreak,
        )
        from reportlab.lib import colors

        kpis = self._ds.kpis()
        enh  = self._ds.enhanced_kpis()
        biz  = enh.business_metrics
        pl   = self._rs.priority_list(top_n=10)
        now  = datetime.now().strftime("%B %d, %Y")
        risk = _risk_label(kpis.churn_rate_pct)

        filename = f"mckinsey_retention_report_{_ts()}.pdf"
        filepath = os.path.join(_ensure_reports_dir(), filename)

        W, H = A4
        doc = SimpleDocTemplate(
            filepath, pagesize=A4,
            leftMargin=2*cm, rightMargin=2*cm,
            topMargin=2*cm, bottomMargin=2*cm,
        )

        def hex2rgb(h):
            h = h.lstrip("#")
            return colors.Color(*[int(h[i:i+2], 16)/255 for i in (0, 2, 4)])

        navy  = hex2rgb(NAVY)
        teal  = hex2rgb(TEAL)
        gold  = hex2rgb(GOLD)
        light = hex2rgb(LIGHT)
        grey  = hex2rgb(GREY)
        red_c = hex2rgb(RED)
        grn_c = hex2rgb(GREEN)

        styles = getSampleStyleSheet()

        def sty(name, **kw):
            return ParagraphStyle(name, parent=styles["Normal"], **kw)

        S = {
            "cover_title": sty("ct", fontSize=28, textColor=navy, leading=34, fontName="Helvetica-Bold"),
            "cover_sub":   sty("cs", fontSize=13, textColor=teal, leading=18, fontName="Helvetica"),
            "cover_meta":  sty("cm", fontSize=9,  textColor=colors.HexColor(MUTED), leading=13),
            "section":     sty("sec", fontSize=13, textColor=navy, fontName="Helvetica-Bold", spaceBefore=14, spaceAfter=4),
            "sub":         sty("sub", fontSize=10, textColor=teal, fontName="Helvetica-Bold", spaceBefore=6, spaceAfter=2),
            "body":        sty("bd",  fontSize=9,  textColor=hex2rgb(TEXT), leading=14, spaceAfter=4),
            "bullet":      sty("bl",  fontSize=9,  textColor=hex2rgb(TEXT), leading=14, leftIndent=12, bulletIndent=0, spaceAfter=2),
            "caption":     sty("cap", fontSize=7.5, textColor=hex2rgb(MUTED), alignment=TA_RIGHT),
            "center":      sty("ctr", fontSize=9,  alignment=TA_CENTER),
            "risk_critical": sty("rc", fontSize=11, textColor=red_c, fontName="Helvetica-Bold"),
            "risk_moderate": sty("rm", fontSize=11, textColor=grn_c, fontName="Helvetica-Bold"),
        }

        def tbl_style(header_color=navy, alt=light):
            return TableStyle([
                ("BACKGROUND",    (0,0), (-1,0),  header_color),
                ("TEXTCOLOR",     (0,0), (-1,0),  colors.white),
                ("FONTNAME",      (0,0), (-1,0),  "Helvetica-Bold"),
                ("FONTSIZE",      (0,0), (-1,0),  9),
                ("ROWBACKGROUNDS",(0,1), (-1,-1),  [colors.white, alt]),
                ("FONTSIZE",      (0,1), (-1,-1),  8.5),
                ("TOPPADDING",    (0,0), (-1,-1),  5),
                ("BOTTOMPADDING", (0,0), (-1,-1),  5),
                ("LEFTPADDING",   (0,0), (-1,-1),  7),
                ("RIGHTPADDING",  (0,0), (-1,-1),  7),
                ("GRID",          (0,0), (-1,-1),  0.4, colors.HexColor("#D1D5DB")),
                ("ROWHEIGHT",     (0,0), (-1,-1),  16),
            ])

        story = []
        PB = PageBreak()
        HR = lambda: HRFlowable(width="100%", thickness=0.5, color=teal, spaceAfter=8)
        SP = lambda n=0.3: Spacer(1, n*cm)

        # ── COVER ──────────────────────────────────────────────────────────────
        story += [
            SP(3),
            Paragraph(req.title or "Customer Retention<br/>Strategic Intelligence Report", S["cover_title"]),
            SP(0.4),
            Paragraph("Executive Briefing — Confidential", S["cover_sub"]),
            SP(0.5),
            HR(),
            SP(0.3),
            Paragraph(f"Prepared by: Customer Retention AI Agent", S["cover_meta"]),
            Paragraph(f"Date: {now}", S["cover_meta"]),
            Paragraph(f"Dataset: {kpis.total_customers:,} customers analysed", S["cover_meta"]),
            Paragraph(f"Risk Classification: {risk}", S["cover_meta"]),
            SP(2),
            Paragraph("STRICTLY CONFIDENTIAL — FOR INTERNAL USE ONLY", S["caption"]),
            PB,
        ]

        # ── 1. EXECUTIVE SUMMARY ───────────────────────────────────────────────
        story += [
            Paragraph("1. Executive Summary", S["section"]),
            HR(),
            Paragraph(
                f"This report presents a comprehensive analysis of customer churn across a "
                f"{kpis.total_customers:,}-customer portfolio. Current churn rate stands at "
                f"<b>{kpis.churn_rate_pct:.1f}%</b>, classified as <b>{risk}</b>. "
                f"Revenue at risk from confirmed churned accounts is "
                f"<b>${kpis.revenue_at_risk:,.0f}</b>, with an additional "
                f"<b>${biz.revenue_at_risk:,.0f}</b> in exposure from the top-quartile "
                f"risk cohort.", S["body"]),
            SP(0.2),
            Paragraph(
                f"The highest churn concentration is in the "
                f"<b>{biz.churn_concentration.top_region}</b> region "
                f"({biz.churn_concentration.pct_of_total_churned:.1f}% of total churned). "
                f"{pl.high_priority_count} customers require immediate retention outreach. "
                f"Executing the recommended actions in this report is projected to recover "
                f"<b>${kpis.revenue_at_risk * 0.20:,.0f} – ${kpis.revenue_at_risk * 0.35:,.0f}</b> "
                f"in annual revenue.", S["body"]),
        ]

        # Risk box
        risk_data = [
            ["Risk Classification", "Churn Rate", "Revenue at Risk", "Priority Customers"],
            [risk,
             f"{kpis.churn_rate_pct:.1f}%",
             f"${kpis.revenue_at_risk:,.0f}",
             f"{pl.high_priority_count} High / {pl.medium_priority_count} Medium"],
        ]
        rt = Table(risk_data, colWidths=[(W-4*cm)/4]*4)
        rt.setStyle(tbl_style(header_color=navy))
        story += [SP(0.3), rt, SP(0.5)]

        # ── 2. KPI DASHBOARD ───────────────────────────────────────────────────
        story += [
            Paragraph("2. KPI Dashboard", S["section"]),
            HR(),
        ]
        kpi_data = [
            ["Metric",                  "Current Value",     "Benchmark",    "Status"],
            ["Total Customers",          f"{kpis.total_customers:,}",          "—",       "✓"],
            ["Churn Rate",               f"{kpis.churn_rate_pct:.2f}%",        "< 15%",  "⚠" if kpis.churn_rate_pct >= 15 else "✓"],
            ["Average CLV",              f"${kpis.avg_clv:,.0f}",              "—",       "✓"],
            ["Revenue at Risk (Churned)",f"${kpis.revenue_at_risk:,.0f}",      "< 5% ARR","⚠"],
            ["Top-Quartile Risk Revenue",f"${biz.revenue_at_risk:,.0f}",       "< 10%",  "⚠"],
            ["Avg Churn Risk Score",     f"{kpis.avg_churn_risk_score:.3f}",   "< 0.30", "⚠" if kpis.avg_churn_risk_score > 0.3 else "✓"],
            ["Avg Satisfaction Score",   f"{kpis.avg_satisfaction_score:.1f}/10","≥ 7.0","⚠" if kpis.avg_satisfaction_score < 7 else "✓"],
            ["High-Value Customers",     f"{biz.high_value_customers:,} ({biz.high_value_pct:.1f}%)","—","✓"],
        ]
        kt = Table(kpi_data, colWidths=[(W-4*cm)*x for x in [0.38, 0.22, 0.22, 0.18]])
        kt.setStyle(tbl_style())
        story += [kt, SP(0.5)]

        # ── 3. CHURN ANALYSIS ─────────────────────────────────────────────────
        story += [
            Paragraph("3. Churn Analysis by Segment & Region", S["section"]),
            HR(),
        ]

        seg_items = sorted(kpis.churn_by_segment.items(), key=lambda x: -x[1])
        reg_items = sorted(kpis.churn_by_region.items(),  key=lambda x: -x[1])

        seg_data = [["Segment", "Churn Rate", "Risk"]] + [
            [s, f"{v*100:.1f}%", _risk_label(v*100)] for s, v in seg_items
        ]
        reg_data = [["Region", "Churn Rate", "Risk"]] + [
            [r, f"{v*100:.1f}%", _risk_label(v*100)] for r, v in reg_items
        ]

        half_w = (W - 4*cm - 0.5*cm) / 2
        seg_t = Table(seg_data, colWidths=[half_w*0.5, half_w*0.3, half_w*0.2])
        reg_t = Table(reg_data, colWidths=[half_w*0.5, half_w*0.3, half_w*0.2])
        seg_t.setStyle(tbl_style(header_color=teal))
        reg_t.setStyle(tbl_style(header_color=teal))

        combined = Table([[seg_t, Spacer(0.5*cm, 1), reg_t]])
        combined.setStyle(TableStyle([("VALIGN", (0,0), (-1,-1), "TOP")]))
        story += [combined, SP(0.5)]

        # ── 4. FINANCIAL IMPACT ────────────────────────────────────────────────
        story += [
            Paragraph("4. Financial Impact Analysis", S["section"]),
            HR(),
        ]
        recovery_low  = kpis.revenue_at_risk * 0.20
        recovery_high = kpis.revenue_at_risk * 0.35

        fin_data = [
            ["Scenario",                       "Revenue Impact"],
            ["Confirmed Churn (Churned CLV)",  f"${kpis.revenue_at_risk:,.0f}"],
            ["High-Risk Cohort Exposure",       f"${biz.revenue_at_risk:,.0f}"],
            ["Total Revenue Exposure",          f"${kpis.revenue_at_risk + biz.revenue_at_risk:,.0f}"],
            ["Conservative Recovery (20%)",    f"${recovery_low:,.0f}"],
            ["Optimistic Recovery (35%)",       f"${recovery_high:,.0f}"],
        ]
        ft = Table(fin_data, colWidths=[(W-4*cm)*0.65, (W-4*cm)*0.35])
        ft.setStyle(tbl_style())
        story += [ft, SP(0.5)]

        # ── 5. PRIORITY CUSTOMER LIST ──────────────────────────────────────────
        story += [
            Paragraph("5. Priority Customer Outreach List (Top 10)", S["section"]),
            HR(),
        ]
        cust_data = [["#", "Customer ID", "Segment", "Region", "CLV", "Risk Score", "Action"]]
        for i, c in enumerate(pl.customers[:10], 1):
            cust_data.append([
                str(i), c.customer_id, c.customer_segment, c.region,
                f"${c.estimated_clv:,.0f}", f"{c.churn_risk_score:.2f}",
                c.recommended_action[:30] + "…" if len(c.recommended_action) > 30 else c.recommended_action,
            ])
        ct = Table(cust_data, colWidths=[(W-4*cm)*x for x in [0.05, 0.14, 0.14, 0.12, 0.12, 0.1, 0.33]])
        ct.setStyle(tbl_style())
        story += [ct, SP(0.5)]

        # ── 6. STRATEGIC RECOMMENDATIONS ──────────────────────────────────────
        story += [
            Paragraph("6. Strategic Recommendations", S["section"]),
            HR(),
        ]
        recommendations = [
            ("Immediate — 0–30 Days",
             f"Deploy personalised Premium Retention Offers to all {pl.high_priority_count} "
             f"High Priority customers. Target CLV > $500 accounts first. "
             f"Expected recovery: ${kpis.revenue_at_risk * 0.12:,.0f}."),
            ("Short-Term — 30–90 Days",
             f"Launch regional retention campaign in {biz.churn_concentration.top_region} "
             f"— the highest-concentration churn zone. "
             f"Deploy satisfaction-improvement programme for customers scoring < 7.0."),
            ("Medium-Term — 90–180 Days",
             "Implement predictive churn intervention: weekly risk-score refresh, "
             "automated outreach triggers at risk > 0.70. "
             "Build loyalty tier programme for top-25% CLV customers."),
            ("Long-Term — 180+ Days",
             "Redesign contract renewal process to reduce monthly plan churn. "
             "Build NPS feedback loop feeding into product roadmap. "
             "Target churn rate reduction to < 15% within 12 months."),
        ]
        for title, body in recommendations:
            story += [
                Paragraph(title, S["sub"]),
                Paragraph(body, S["body"]),
            ]

        story += [
            SP(0.5),
            HR(),
            Paragraph(
                f"Report generated: {now} · Customer Retention AI Agent · CONFIDENTIAL",
                S["caption"],
            ),
        ]

        doc.build(story)
        return ReportResponse(
            report_type="pdf",
            filename=filename,
            download_path=f"/reports/export/{filename}",
            generated_at=datetime.now(timezone.utc).isoformat(),
            size_bytes=os.path.getsize(filepath),
        )

    # ──────────────────────────────────────────────────────────────────────────
    # Excel — multi-sheet McKinsey workbook
    # ──────────────────────────────────────────────────────────────────────────

    def generate_csv(self, req: ReportRequest) -> ReportResponse:
        """Generates a rich multi-sheet Excel workbook (not just CSV)."""
        import openpyxl
        from openpyxl.styles import (
            Font, PatternFill, Alignment, Border, Side, numbers,
        )
        from openpyxl.utils import get_column_letter
        from openpyxl.chart import BarChart, Reference
        from openpyxl.chart.series import DataPoint

        kpis = self._ds.kpis()
        enh  = self._ds.enhanced_kpis()
        biz  = enh.business_metrics
        pl   = self._rs.priority_list(top_n=50)
        now  = datetime.now().strftime("%B %d, %Y")
        risk = _risk_label(kpis.churn_rate_pct)

        wb = openpyxl.Workbook()

        # Shared styles
        def fill(hex_color):
            return PatternFill("solid", fgColor=hex_color.lstrip("#"))

        def font(bold=False, color="1A1A2E", size=10, italic=False):
            return Font(bold=bold, color=color.lstrip("#"), size=size, italic=italic)

        def border():
            s = Side(style="thin", color="D1D5DB")
            return Border(left=s, right=s, top=s, bottom=s)

        def center():
            return Alignment(horizontal="center", vertical="center", wrap_text=True)

        def left():
            return Alignment(horizontal="left", vertical="center", wrap_text=True)

        def write_header(ws, row, col, text, width=None):
            c = ws.cell(row=row, column=col, value=text)
            c.font = font(bold=True, color=WHITE, size=10)
            c.fill = fill(NAVY)
            c.alignment = center()
            c.border = border()
            if width and col <= ws.max_column:
                ws.column_dimensions[get_column_letter(col)].width = width
            return c

        def write_cell(ws, row, col, value, bold=False, color=TEXT, bg=None, align="left", num_fmt=None):
            c = ws.cell(row=row, column=col, value=value)
            c.font = font(bold=bold, color=color, size=9)
            c.alignment = Alignment(horizontal=align, vertical="center")
            c.border = border()
            if bg:
                c.fill = fill(bg)
            if num_fmt:
                c.number_format = num_fmt
            return c

        def set_col_width(ws, col, width):
            ws.column_dimensions[get_column_letter(col)].width = width

        def freeze(ws, cell="A2"):
            ws.freeze_panes = cell

        # ── Sheet 1: Executive Summary ─────────────────────────────────────────
        ws1 = wb.active
        ws1.title = "Executive Summary"
        ws1.sheet_view.showGridLines = False

        # Title block
        ws1.merge_cells("A1:F1")
        t = ws1["A1"]
        t.value = "CUSTOMER RETENTION STRATEGIC INTELLIGENCE REPORT"
        t.font = Font(bold=True, size=16, color=WHITE.lstrip("#"))
        t.fill = fill(NAVY)
        t.alignment = center()
        ws1.row_dimensions[1].height = 36

        ws1.merge_cells("A2:F2")
        ws1["A2"].value = f"Prepared by Customer Retention AI Agent  ·  {now}  ·  CONFIDENTIAL"
        ws1["A2"].font = font(italic=True, color=MUTED, size=9)
        ws1["A2"].alignment = center()
        ws1["A2"].fill = fill(LIGHT.lstrip("#"))

        row = 4
        headers = ["Metric", "Value", "Benchmark", "Status", "Risk Level", "Priority"]
        for c, h in enumerate(headers, 1):
            write_header(ws1, row, c, h)
        set_col_width(ws1, 1, 32); set_col_width(ws1, 2, 18)
        set_col_width(ws1, 3, 18); set_col_width(ws1, 4, 12)
        set_col_width(ws1, 5, 16); set_col_width(ws1, 6, 16)

        kpi_rows = [
            ("Total Customers",          f"{kpis.total_customers:,}",          "—",         "✓", "—",      "—"),
            ("Churn Rate",               f"{kpis.churn_rate_pct:.2f}%",        "< 15.0%",  "⚠", risk,     "IMMEDIATE"),
            ("Revenue at Risk (Churned)",f"${kpis.revenue_at_risk:,.0f}",      "< 5% ARR", "⚠", "HIGH",   "IMMEDIATE"),
            ("Top-Quartile Risk Exposure",f"${biz.revenue_at_risk:,.0f}",      "< 10% ARR","⚠", "HIGH",   "SHORT-TERM"),
            ("Average CLV",              f"${kpis.avg_clv:,.0f}",              "—",         "✓", "—",      "MONITOR"),
            ("Avg Churn Risk Score",     f"{kpis.avg_churn_risk_score:.3f}",   "< 0.30",   "⚠", "ELEVATED","SHORT-TERM"),
            ("Avg Satisfaction Score",   f"{kpis.avg_satisfaction_score:.1f}", "≥ 7.0",    "⚠", "ELEVATED","MEDIUM-TERM"),
            ("High-Value Customers",     f"{biz.high_value_customers:,} ({biz.high_value_pct:.1f}%)","—","✓","—","PROTECT"),
            ("High Priority Outreach",   f"{pl.high_priority_count} customers","—",         "⚠", "HIGH",   "IMMEDIATE"),
            ("Medium Priority Outreach", f"{pl.medium_priority_count} customers","—",       "~", "MODERATE","SHORT-TERM"),
        ]
        for i, r in enumerate(kpi_rows):
            dr = row + 1 + i
            bg = GREY if i % 2 == 0 else WHITE
            for c, val in enumerate(r, 1):
                clr = RED if val in ("⚠","IMMEDIATE","HIGH","CRITICAL") else (
                      GREEN if val in ("✓","—") else TEXT)
                write_cell(ws1, dr, c, val, color=clr, bg=bg,
                           align="center" if c > 1 else "left")

        # Recovery projection
        row2 = row + len(kpi_rows) + 3
        ws1.merge_cells(f"A{row2}:F{row2}")
        ws1[f"A{row2}"].value = "RECOVERY PROJECTION"
        ws1[f"A{row2}"].font = font(bold=True, color=WHITE, size=11)
        ws1[f"A{row2}"].fill = fill(TEAL)
        ws1[f"A{row2}"].alignment = center()

        proj_rows = [
            ("Conservative Recovery (20%)", f"${kpis.revenue_at_risk * 0.20:,.0f}",
             "Immediate outreach to High Priority queue"),
            ("Base Case Recovery (28%)",    f"${kpis.revenue_at_risk * 0.28:,.0f}",
             "Full campaign + satisfaction improvement"),
            ("Optimistic Recovery (35%)",   f"${kpis.revenue_at_risk * 0.35:,.0f}",
             "Full programme + predictive intervention"),
        ]
        write_header(ws1, row2+1, 1, "Scenario")
        write_header(ws1, row2+1, 2, "Revenue Protected")
        write_header(ws1, row2+1, 3, "Conditions")
        ws1.merge_cells(f"C{row2+1}:F{row2+1}")
        for i, (sc, val, cond) in enumerate(proj_rows):
            write_cell(ws1, row2+2+i, 1, sc, bold=True)
            write_cell(ws1, row2+2+i, 2, val, color=GREEN, bold=True, align="center")
            write_cell(ws1, row2+2+i, 3, cond)
            ws1.merge_cells(f"C{row2+2+i}:F{row2+2+i}")

        freeze(ws1, "A3")

        # ── Sheet 2: Churn Breakdown ───────────────────────────────────────────
        ws2 = wb.create_sheet("Churn Breakdown")
        ws2.sheet_view.showGridLines = False

        ws2.merge_cells("A1:E1")
        ws2["A1"].value = "Churn Rate Analysis — By Segment & Region"
        ws2["A1"].font = font(bold=True, color=WHITE, size=13)
        ws2["A1"].fill = fill(NAVY)
        ws2["A1"].alignment = center()
        ws2.row_dimensions[1].height = 28

        # Segment table
        ws2["A3"].value = "BY CUSTOMER SEGMENT"
        ws2["A3"].font = font(bold=True, color=TEAL, size=10)
        for c, h in enumerate(["Segment", "Churn Rate", "Risk Level", "Churned Est.", "Action"], 1):
            write_header(ws2, 4, c, h)
        set_col_width(ws2, 1, 24); set_col_width(ws2, 2, 16)
        set_col_width(ws2, 3, 16); set_col_width(ws2, 4, 16); set_col_width(ws2, 5, 28)

        seg_sorted = sorted(kpis.churn_by_segment.items(), key=lambda x: -x[1])
        for i, (seg, rate) in enumerate(seg_sorted):
            r = 5 + i
            est_churned = int(kpis.total_customers * rate * (1/len(seg_sorted)))
            action = ("Premium Retention Offer" if rate > 0.25 else
                      "Proactive Outreach" if rate > 0.15 else "Monitor")
            bg = GREY if i % 2 == 0 else WHITE
            write_cell(ws2, r, 1, seg, bg=bg)
            write_cell(ws2, r, 2, f"{rate*100:.1f}%", align="center", bg=bg,
                       color=RED if rate > 0.20 else TEXT)
            write_cell(ws2, r, 3, _risk_label(rate*100), align="center", bg=bg,
                       color=RED if rate > 0.20 else TEXT)
            write_cell(ws2, r, 4, est_churned, align="center", bg=bg)
            write_cell(ws2, r, 5, action, bg=bg)

        # Region table
        row_r = 5 + len(seg_sorted) + 3
        ws2[f"A{row_r-1}"].value = "BY REGION"
        ws2[f"A{row_r-1}"].font = font(bold=True, color=TEAL, size=10)
        for c, h in enumerate(["Region", "Churn Rate", "Risk Level", "Churned Est.", "Campaign Priority"], 1):
            write_header(ws2, row_r, c, h)

        reg_sorted = sorted(kpis.churn_by_region.items(), key=lambda x: -x[1])
        for i, (reg, rate) in enumerate(reg_sorted):
            r = row_r + 1 + i
            est_churned = int(kpis.total_customers * rate * (1/len(reg_sorted)))
            pri = "P1 — Launch Now" if rate > 0.22 else "P2 — Next Cycle" if rate > 0.15 else "P3 — Monitor"
            bg = GREY if i % 2 == 0 else WHITE
            write_cell(ws2, r, 1, reg, bg=bg)
            write_cell(ws2, r, 2, f"{rate*100:.1f}%", align="center", bg=bg,
                       color=RED if rate > 0.20 else TEXT)
            write_cell(ws2, r, 3, _risk_label(rate*100), align="center", bg=bg)
            write_cell(ws2, r, 4, est_churned, align="center", bg=bg)
            write_cell(ws2, r, 5, pri, bg=bg,
                       color=RED if "P1" in pri else (TEAL if "P2" in pri else TEXT))

        freeze(ws2, "A2")

        # ── Sheet 3: Priority Customer List ────────────────────────────────────
        ws3 = wb.create_sheet("Priority Customer List")
        ws3.sheet_view.showGridLines = False

        ws3.merge_cells("A1:H1")
        ws3["A1"].value = f"Priority Customer Outreach List — Top {len(pl.customers)} Accounts"
        ws3["A1"].font = font(bold=True, color=WHITE, size=13)
        ws3["A1"].fill = fill(NAVY)
        ws3["A1"].alignment = center()
        ws3.row_dimensions[1].height = 28

        hdrs = ["#", "Customer ID", "Segment", "Plan Type", "Region",
                "CLV ($)", "Risk Score", "Priority", "Recommended Action"]
        widths = [5, 18, 16, 14, 14, 14, 12, 12, 36]
        for c, (h, w) in enumerate(zip(hdrs, widths), 1):
            write_header(ws3, 2, c, h)
            set_col_width(ws3, c, w)

        for i, cust in enumerate(pl.customers):
            r = 3 + i
            bg = GREY if i % 2 == 0 else WHITE
            pri_color = RED if cust.priority_class == "High" else (GOLD if cust.priority_class == "Medium" else GREEN)
            write_cell(ws3, r, 1, i+1, align="center", bg=bg)
            write_cell(ws3, r, 2, cust.customer_id, bg=bg)
            write_cell(ws3, r, 3, cust.customer_segment, bg=bg)
            write_cell(ws3, r, 4, cust.plan_type, bg=bg)
            write_cell(ws3, r, 5, cust.region, bg=bg)
            write_cell(ws3, r, 6, cust.estimated_clv, align="center", bg=bg, num_fmt='"$"#,##0')
            write_cell(ws3, r, 7, cust.churn_risk_score, align="center", bg=bg, num_fmt="0.000")
            write_cell(ws3, r, 8, cust.priority_class, align="center", bg=bg,
                       bold=True, color=pri_color)
            write_cell(ws3, r, 9, cust.recommended_action, bg=bg)

        freeze(ws3, "A3")

        # ── Sheet 4: Financial Model ────────────────────────────────────────────
        ws4 = wb.create_sheet("Financial Model")
        ws4.sheet_view.showGridLines = False

        ws4.merge_cells("A1:D1")
        ws4["A1"].value = "Revenue Impact & Recovery Model"
        ws4["A1"].font = font(bold=True, color=WHITE, size=13)
        ws4["A1"].fill = fill(NAVY)
        ws4["A1"].alignment = center()
        ws4.row_dimensions[1].height = 28

        set_col_width(ws4, 1, 34); set_col_width(ws4, 2, 22)
        set_col_width(ws4, 3, 22); set_col_width(ws4, 4, 22)

        # Revenue waterfall data
        fin_sections = [
            ("CURRENT STATE", None),
            ("Total Customer Base Value", kpis.avg_clv * kpis.total_customers, None),
            ("Revenue Lost (Churned Customers)", -kpis.revenue_at_risk, None),
            ("At-Risk Cohort Exposure", -biz.revenue_at_risk, None),
            ("Net Protected Revenue", kpis.avg_clv * kpis.total_customers - kpis.revenue_at_risk - biz.revenue_at_risk, None),
            (None, None, None),
            ("INTERVENTION SCENARIOS", None, None),
            ("Budget: $250,000 Campaign (10% disc.)", None, None),
            ("  Customers Reachable", int(250000 / max(kpis.avg_clv * 0.10, 1)), None),
            ("  Expected Retained (18% success)", int(int(250000 / max(kpis.avg_clv * 0.10, 1)) * 0.18), None),
            ("  Revenue Protected",  int(250000 / max(kpis.avg_clv * 0.10, 1)) * 0.18 * kpis.avg_clv, None),
            ("  Estimated ROI", (int(250000 / max(kpis.avg_clv * 0.10, 1)) * 0.18 * kpis.avg_clv) / 250000, None),
            (None, None, None),
            ("Budget: $500,000 Campaign (10% disc.)", None, None),
            ("  Customers Reachable", int(500000 / max(kpis.avg_clv * 0.10, 1)), None),
            ("  Expected Retained (18% success)", int(int(500000 / max(kpis.avg_clv * 0.10, 1)) * 0.18), None),
            ("  Revenue Protected",  int(500000 / max(kpis.avg_clv * 0.10, 1)) * 0.18 * kpis.avg_clv, None),
            ("  Estimated ROI", (int(500000 / max(kpis.avg_clv * 0.10, 1)) * 0.18 * kpis.avg_clv) / 500000, None),
        ]

        for c, h in enumerate(["Item", "Value", "", ""], 1):
            write_header(ws4, 2, c, h)

        for i, row_data in enumerate(fin_sections):
            r = 3 + i
            label = row_data[0]
            val   = row_data[1]
            bg = GREY if i % 2 == 0 else WHITE
            if label is None:
                ws4.row_dimensions[r].height = 6
                continue
            is_section = not label.startswith(" ") and val is None
            c = ws4.cell(row=r, column=1, value=label)
            c.font = font(bold=is_section, color=TEAL if is_section else TEXT, size=9)
            c.fill = fill(LIGHT.lstrip("#")) if is_section else fill(bg.lstrip("#"))
            c.alignment = left()
            c.border = border()
            if val is not None:
                v = ws4.cell(row=r, column=2, value=val)
                v.font = font(size=9, color=RED if isinstance(val, (int, float)) and val < 0 else TEXT)
                v.alignment = Alignment(horizontal="right", vertical="center")
                v.border = border()
                v.fill = fill(bg.lstrip("#"))
                if isinstance(val, float) and val < 100:
                    v.number_format = '0.00"×"'
                else:
                    v.number_format = '"$"#,##0'
            ws4.merge_cells(f"B{r}:D{r}")

        freeze(ws4, "A3")

        # ── Sheet 5: Action Plan ───────────────────────────────────────────────
        ws5 = wb.create_sheet("90-Day Action Plan")
        ws5.sheet_view.showGridLines = False

        ws5.merge_cells("A1:F1")
        ws5["A1"].value = "90-Day Customer Retention Action Plan"
        ws5["A1"].font = font(bold=True, color=WHITE, size=13)
        ws5["A1"].fill = fill(NAVY)
        ws5["A1"].alignment = center()
        ws5.row_dimensions[1].height = 28

        plan_hdrs = ["Timeframe", "Priority", "Initiative", "Target Segment",
                     "KPI Impact", "Owner"]
        plan_widths = [18, 12, 40, 20, 24, 18]
        for c, (h, w) in enumerate(zip(plan_hdrs, plan_widths), 1):
            write_header(ws5, 2, c, h)
            set_col_width(ws5, c, w)

        top_seg = max(kpis.churn_by_segment.items(), key=lambda x: x[1], default=("All", 0))
        top_reg = max(kpis.churn_by_region.items(),  key=lambda x: x[1], default=("All", 0))

        action_plan = [
            ("Days 1–14",   "CRITICAL", "Immediate outreach to High Priority queue",
             "High-risk accounts", f"Protect ${kpis.revenue_at_risk*0.08:,.0f}", "Retention Team"),
            ("Days 1–30",   "HIGH",     "Premium Retention Offer deployment",
             f"{top_seg[0]} segment", f"Reduce churn by 3–5%", "CRM / Marketing"),
            ("Days 15–45",  "HIGH",     f"Regional campaign — {top_reg[0]}",
             f"{top_reg[0]} region", "Reduce regional churn 8%", "Regional Mgr"),
            ("Days 30–60",  "HIGH",     "Satisfaction Recovery Programme",
             "Score < 7.0", "+1.5 pts avg satisfaction", "CX Team"),
            ("Days 30–60",  "MEDIUM",   "Payment Support Plan rollout",
             "Failed payments > 3", "Reduce payment churn 40%", "Finance / CS"),
            ("Days 45–75",  "MEDIUM",   "Predictive churn model weekly refresh",
             "All customers", "Catch risk before it spikes", "Data / Analytics"),
            ("Days 60–90",  "MEDIUM",   "Upsell campaign for low-risk high-CLV",
             "Standard Plan, risk < 0.3", "+12% upsell conversion", "Sales"),
            ("Days 75–90",  "LOW",      "NPS survey + product feedback loop",
             "All segments", "Identify CX improvement areas", "Product"),
            ("Ongoing",     "LOW",      "Monthly exec dashboard + KPI review",
             "Leadership", "Churn < 15% in 12 months", "CEO / COO"),
        ]

        for i, row_data in enumerate(action_plan):
            r = 3 + i
            bg = GREY if i % 2 == 0 else WHITE
            for c, val in enumerate(row_data, 1):
                clr = RED if val == "CRITICAL" else (GOLD if val == "HIGH" else (TEAL if val == "MEDIUM" else MUTED))
                is_pri = (c == 2)
                write_cell(ws5, r, c, val, bg=bg,
                           color=clr if is_pri else TEXT,
                           bold=is_pri, align="center" if is_pri else "left")

        freeze(ws5, "A3")

        # Save
        filename = f"mckinsey_retention_workbook_{_ts()}.xlsx"
        filepath = os.path.join(_ensure_reports_dir(), filename)
        wb.save(filepath)

        return ReportResponse(
            report_type="xlsx",
            filename=filename,
            download_path=f"/reports/export/{filename}",
            generated_at=datetime.now(timezone.utc).isoformat(),
            size_bytes=os.path.getsize(filepath),
        )

    # ──────────────────────────────────────────────────────────────────────────
    # PowerPoint — McKinsey-style 10-slide deck
    # ──────────────────────────────────────────────────────────────────────────

    def generate_ppt(self, req: ReportRequest) -> ReportResponse:
        try:
            return self._generate_pptx(req)
        except ImportError:
            return self._generate_markdown(req)

    def _generate_pptx(self, req: ReportRequest) -> ReportResponse:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        def rgb(h):
            h = h.lstrip("#")
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        kpis = self._ds.kpis()
        enh  = self._ds.enhanced_kpis()
        biz  = enh.business_metrics
        pl   = self._rs.priority_list(top_n=8)
        now  = datetime.now().strftime("%B %d, %Y")
        risk = _risk_label(kpis.churn_rate_pct)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        BLANK = prs.slide_layouts[6]

        def add_slide():
            return prs.slides.add_slide(BLANK)

        def rect(slide, x, y, w, h, fill_hex=None, line_hex=None):
            from pptx.util import Inches
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.line.fill.background()
            if fill_hex:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(fill_hex)
            else:
                shape.fill.background()
            if line_hex:
                shape.line.color.rgb = rgb(line_hex)
                shape.line.width = Pt(0.75)
            else:
                shape.line.fill.background()
            return shape

        def txt(slide, text, x, y, w, h, size=12, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf  = box.text_frame
            tf.word_wrap = wrap
            p   = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.bold   = bold
            run.font.italic = italic
            run.font.size   = Pt(size)
            run.font.color.rgb = rgb(color)
            return box

        def title_bar(slide, text, sub=None):
            rect(slide, 0, 0, 13.33, 1.3, fill_hex=NAVY)
            txt(slide, text, 0.4, 0.12, 12.5, 0.7, size=22, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
            if sub:
                txt(slide, sub, 0.4, 0.78, 12.5, 0.4, size=11, color=LIGHT,
                    align=PP_ALIGN.LEFT)
            rect(slide, 0, 1.28, 13.33, 0.05, fill_hex=GOLD)

        def bullet_box(slide, title, bullets, x, y, w, h, title_color=NAVY):
            rect(slide, x, y, w, 0.38, fill_hex=LIGHT.lstrip("#"))
            txt(slide, title, x+0.1, y+0.05, w-0.2, 0.28, size=10, bold=True, color=title_color)
            for i, b in enumerate(bullets):
                txt(slide, f"▸  {b}", x+0.1, y+0.45+i*0.38, w-0.2, 0.35, size=9, color=TEXT)

        def kpi_box(slide, label, value, x, y, w=2.8, val_color=NAVY):
            rect(slide, x, y, w, 1.15, fill_hex=LIGHT.lstrip("#"), line_hex=TEAL)
            txt(slide, label, x+0.1, y+0.08, w-0.2, 0.3, size=8.5, color=MUTED)
            txt(slide, value, x+0.1, y+0.38, w-0.2, 0.65, size=18, bold=True,
                color=val_color, align=PP_ALIGN.LEFT)

        # ── Slide 1: Cover ─────────────────────────────────────────────────────
        sl = add_slide()
        rect(sl, 0, 0, 13.33, 7.5, fill_hex=NAVY)
        rect(sl, 0, 5.8, 13.33, 1.7, fill_hex=TEAL)
        txt(sl, "CUSTOMER RETENTION", 1, 1.2, 11, 1, size=36, bold=True, color=WHITE)
        txt(sl, "Strategic Intelligence Report", 1, 2.2, 11, 0.7, size=24, color=LIGHT)
        rect(sl, 1, 3.1, 4, 0.06, fill_hex=GOLD)
        txt(sl, f"Executive Briefing  ·  {now}", 1, 3.3, 10, 0.4, size=13, color=LIGHT)
        txt(sl, f"Portfolio: {kpis.total_customers:,} Customers Analysed", 1, 3.75, 10, 0.4, size=12, color=LIGHT)
        txt(sl, f"Risk Classification: {risk}", 1, 4.15, 10, 0.4, size=12, bold=True,
            color=GOLD if risk in ("HIGH","CRITICAL") else LIGHT)
        txt(sl, "STRICTLY CONFIDENTIAL — PREPARED BY CUSTOMER RETENTION AI AGENT",
            0.5, 6.0, 12, 0.4, size=9, color=WHITE, align=PP_ALIGN.CENTER)

        # ── Slide 2: Agenda ────────────────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Agenda", "Today's Strategic Briefing")
        items = [
            ("01", "Executive Summary & Risk Classification"),
            ("02", "KPI Dashboard — Current State"),
            ("03", "Churn Analysis by Segment & Region"),
            ("04", "Financial Impact & Revenue at Risk"),
            ("05", "Priority Customer Outreach List"),
            ("06", "Strategic Recommendations"),
            ("07", "90-Day Execution Roadmap"),
            ("08", "Expected Outcomes & ROI Model"),
        ]
        for i, (num, label) in enumerate(items):
            col = i // 4
            row_i = i % 4
            x = 0.5 + col * 6.3
            y = 1.6 + row_i * 1.3
            rect(sl, x, y, 5.8, 1.1, fill_hex=LIGHT.lstrip("#"), line_hex=TEAL)
            txt(sl, num, x+0.15, y+0.08, 0.7, 0.8, size=22, bold=True, color=TEAL)
            txt(sl, label, x+0.75, y+0.2, 4.8, 0.7, size=11, color=TEXT)

        # ── Slide 3: Executive Summary ─────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Executive Summary", f"Risk Level: {risk}  ·  {kpis.total_customers:,} Customers Analysed")
        rect(sl, 0.4, 1.5, 12.5, 1.0, fill_hex=RED if risk in ("CRITICAL","HIGH") else TEAL)
        txt(sl, f"⚠  RISK CLASSIFICATION: {risk}  —  Churn at {kpis.churn_rate_pct:.1f}%  —  ${kpis.revenue_at_risk:,.0f} Revenue at Risk",
            0.6, 1.62, 12, 0.7, size=13, bold=True, color=WHITE)
        summary = (
            f"This portfolio of {kpis.total_customers:,} customers is experiencing a churn rate of "
            f"{kpis.churn_rate_pct:.1f}%, classified as {risk}. Revenue at risk from confirmed churned "
            f"accounts totals ${kpis.revenue_at_risk:,.0f}, with an additional ${biz.revenue_at_risk:,.0f} "
            f"in exposure from the high-risk cohort. Immediate intervention is required across "
            f"{pl.high_priority_count} High Priority accounts."
        )
        txt(sl, summary, 0.4, 2.65, 8.5, 2.0, size=11, color=TEXT)

        kpi_box(sl, "Churn Rate", f"{kpis.churn_rate_pct:.1f}%",
                9.2, 1.55, val_color=RED if kpis.churn_rate_pct > 20 else NAVY)
        kpi_box(sl, "Revenue at Risk", f"${kpis.revenue_at_risk/1e6:.2f}M", 9.2, 2.85)
        kpi_box(sl, "High Priority", f"{pl.high_priority_count} accounts",
                9.2, 4.15, val_color=RED)

        # ── Slide 4: KPI Dashboard ─────────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "KPI Dashboard — Current State")

        kpis_display = [
            ("Total Customers",    f"{kpis.total_customers:,}",      NAVY),
            ("Churn Rate",         f"{kpis.churn_rate_pct:.1f}%",     RED),
            ("Avg CLV",            f"${kpis.avg_clv:,.0f}",           NAVY),
            ("Revenue at Risk",    f"${kpis.revenue_at_risk/1e6:.2f}M", RED),
            ("Avg Risk Score",     f"{kpis.avg_churn_risk_score:.3f}", TEAL),
            ("Avg Satisfaction",   f"{kpis.avg_satisfaction_score:.1f}/10", TEAL),
            ("High-Value Customers",f"{biz.high_value_customers:,}",  GREEN),
            ("At-Risk Exposure",   f"${biz.revenue_at_risk/1e6:.2f}M", RED),
        ]
        cols = 4
        for i, (label, val, col) in enumerate(kpis_display):
            cx = 0.4 + (i % cols) * 3.15
            cy = 1.55 + (i // cols) * 1.4
            kpi_box(sl, label, val, cx, cy, w=2.9, val_color=col)

        txt(sl, f"Data as of {now}  ·  Customer Retention AI Agent", 0.4, 6.9, 12.5, 0.4,
            size=8, color=MUTED, align=PP_ALIGN.RIGHT)

        # ── Slide 5: Churn by Segment ──────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Churn Analysis — By Segment & Region",
                  "Sorted by churn rate descending")

        seg_sorted = sorted(kpis.churn_by_segment.items(), key=lambda x: -x[1])
        reg_sorted = sorted(kpis.churn_by_region.items(),  key=lambda x: -x[1])

        # Segment bars
        txt(sl, "CHURN BY SEGMENT", 0.4, 1.55, 5.5, 0.35, size=10, bold=True, color=TEAL)
        max_seg = seg_sorted[0][1] if seg_sorted else 1
        for i, (seg, rate) in enumerate(seg_sorted[:6]):
            y = 2.0 + i * 0.72
            bar_w = (rate / max(max_seg, 0.01)) * 4.5
            clr = RED if rate > 0.22 else (GOLD if rate > 0.16 else GREEN)
            rect(sl, 2.0, y+0.08, bar_w, 0.45, fill_hex=clr)
            txt(sl, seg, 0.4, y+0.1, 1.5, 0.35, size=9, color=TEXT)
            txt(sl, f"{rate*100:.1f}%", 2.1+bar_w, y+0.1, 0.8, 0.35, size=9,
                bold=True, color=clr)

        # Region bars
        txt(sl, "CHURN BY REGION", 7.0, 1.55, 5.8, 0.35, size=10, bold=True, color=TEAL)
        max_reg = reg_sorted[0][1] if reg_sorted else 1
        for i, (reg, rate) in enumerate(reg_sorted[:6]):
            y = 2.0 + i * 0.72
            bar_w = (rate / max(max_reg, 0.01)) * 4.5
            clr = RED if rate > 0.22 else (GOLD if rate > 0.16 else GREEN)
            rect(sl, 8.6, y+0.08, bar_w, 0.45, fill_hex=clr)
            txt(sl, reg, 7.0, y+0.1, 1.5, 0.35, size=9, color=TEXT)
            txt(sl, f"{rate*100:.1f}%", 8.7+bar_w, y+0.1, 0.8, 0.35, size=9,
                bold=True, color=clr)

        rect(sl, 6.5, 1.5, 0.04, 5.8, fill_hex=LIGHT.lstrip("#"))

        # ── Slide 6: Financial Impact ──────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Financial Impact Analysis", "Revenue at Risk Waterfall & Recovery Scenarios")

        items_fin = [
            ("Total Portfolio Value",       kpis.avg_clv * kpis.total_customers, TEAL,  True),
            ("Less: Churned CLV",           -kpis.revenue_at_risk,               RED,   False),
            ("Less: High-Risk Cohort",      -biz.revenue_at_risk,                RED,   False),
            ("= Net Protected Value",
             kpis.avg_clv * kpis.total_customers - kpis.revenue_at_risk - biz.revenue_at_risk,
             NAVY, True),
        ]
        for i, (label, val, clr, bold) in enumerate(items_fin):
            y = 1.6 + i * 0.85
            rect(sl, 0.4, y, 6.5, 0.72, fill_hex=LIGHT.lstrip("#") if bold else WHITE,
                 line_hex=TEAL if bold else "D1D5DB")
            txt(sl, label, 0.6, y+0.12, 3.8, 0.5, size=10, bold=bold, color=TEXT)
            txt(sl, f"${abs(val):,.0f}" + (" ▼" if val < 0 else ""),
                4.6, y+0.1, 2.1, 0.5, size=14, bold=True, color=clr,
                align=PP_ALIGN.RIGHT)

        txt(sl, "RECOVERY SCENARIOS", 7.5, 1.55, 5.4, 0.35, size=10, bold=True, color=TEAL)
        scenarios = [
            ("Conservative (20%)", kpis.revenue_at_risk * 0.20, "Immediate outreach only"),
            ("Base Case (28%)",    kpis.revenue_at_risk * 0.28, "Full campaign + CX"),
            ("Optimistic (35%)",   kpis.revenue_at_risk * 0.35, "Full programme + AI"),
        ]
        for i, (label, val, note) in enumerate(scenarios):
            y = 2.0 + i * 1.55
            rect(sl, 7.5, y, 5.4, 1.3, fill_hex=GREEN if i==2 else LIGHT.lstrip("#"),
                 line_hex=GREEN)
            txt(sl, label, 7.65, y+0.08, 5.0, 0.35, size=9.5, bold=True,
                color=WHITE if i==2 else NAVY)
            txt(sl, f"${val:,.0f}", 7.65, y+0.45, 5.0, 0.5, size=18, bold=True,
                color=WHITE if i==2 else GREEN)
            txt(sl, note, 7.65, y+0.95, 5.0, 0.28, size=8, color=WHITE if i==2 else MUTED)

        # ── Slide 7: Priority Customers ────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, f"Priority Customer Outreach — Top {min(8, len(pl.customers))} Accounts",
                  "Sorted by composite priority score (churn risk × CLV × complaints)")

        hdrs_p = ["#", "Customer ID", "Segment", "Region", "CLV", "Risk", "Priority", "Action"]
        col_x  = [0.25, 0.6, 2.2, 3.7, 5.1, 6.4, 7.5, 8.7]
        col_w  = [0.3,  1.5, 1.4, 1.3, 1.2, 1.0, 1.1, 4.3]

        for c, (h, x, w) in enumerate(zip(hdrs_p, col_x, col_w)):
            rect(sl, x, 1.5, w, 0.42, fill_hex=NAVY)
            txt(sl, h, x+0.05, 1.55, w-0.05, 0.32, size=8.5, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

        for i, cust in enumerate(pl.customers[:8]):
            y = 2.0 + i * 0.58
            bg = LIGHT.lstrip("#") if i % 2 == 0 else WHITE
            pri_clr = RED if cust.priority_class=="High" else (GOLD if cust.priority_class=="Medium" else GREEN)
            vals = [str(i+1), cust.customer_id, cust.customer_segment, cust.region,
                    f"${cust.estimated_clv:,.0f}", f"{cust.churn_risk_score:.2f}",
                    cust.priority_class,
                    cust.recommended_action[:40]+"…" if len(cust.recommended_action)>40 else cust.recommended_action]
            for v, x, w in zip(vals, col_x, col_w):
                is_pri = (v == cust.priority_class)
                rect(sl, x, y, w, 0.5, fill_hex=bg, line_hex="D1D5DB")
                txt(sl, v, x+0.04, y+0.08, w-0.06, 0.35, size=8,
                    color=pri_clr if is_pri else TEXT, bold=is_pri,
                    align=PP_ALIGN.CENTER)

        # ── Slide 8: Recommendations ───────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Strategic Recommendations",
                  "Immediate · Short-Term · Medium-Term priorities")

        top_seg = max(kpis.churn_by_segment.items(), key=lambda x: x[1], default=("All",0))
        top_reg = max(kpis.churn_by_region.items(),  key=lambda x: x[1], default=("All",0))

        recs = [
            ("IMMEDIATE — 0–30 Days", RED, [
                f"Deploy Premium Retention Offers to {pl.high_priority_count} High Priority accounts",
                f"Prioritise CLV > ${kpis.avg_clv*1.5:,.0f} accounts first — highest ROI",
                f"Launch regional campaign in {top_reg[0]} (highest churn zone)",
            ]),
            ("SHORT-TERM — 30–90 Days", GOLD, [
                f"Satisfaction Recovery Programme for scores < 7.0",
                f"Payment Support Plans for customers with payment failures",
                f"Predictive risk-score refresh — weekly automated alerts at risk > 0.70",
            ]),
            ("MEDIUM-TERM — 90–180 Days", TEAL, [
                "Loyalty tier programme for top-25% CLV customers",
                "Contract renewal redesign to reduce monthly plan churn",
                "NPS feedback loop integrated into product roadmap",
            ]),
        ]
        for i, (label, clr, bullets) in enumerate(recs):
            x = 0.35 + i * 4.3
            rect(sl, x, 1.5, 4.1, 5.7, fill_hex=LIGHT.lstrip("#"), line_hex=clr)
            rect(sl, x, 1.5, 4.1, 0.45, fill_hex=clr)
            txt(sl, label, x+0.1, 1.55, 3.9, 0.35, size=9, bold=True, color=WHITE)
            for j, b in enumerate(bullets):
                txt(sl, f"▸  {b}", x+0.15, 2.15+j*0.85, 3.8, 0.75, size=9.5, color=TEXT)

        # ── Slide 9: 90-Day Roadmap ────────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "90-Day Execution Roadmap")

        phases = [
            ("Week 1–2",  "LAUNCH", RED,  ["High Priority outreach", "Campaign brief finalised"]),
            ("Week 3–6",  "DEPLOY", GOLD, ["Retention offers live", f"{top_reg[0]} campaign",]),
            ("Week 7–10", "SCALE",  TEAL, ["Sat. recovery programme", "Payment support"]),
            ("Week 11–13","OPTIMISE",GREEN,["Upsell campaign", "Predictive refresh"]),
        ]
        for i, (period, phase, clr, milestones) in enumerate(phases):
            x = 0.4 + i * 3.2
            rect(sl, x, 1.5, 3.0, 5.7, fill_hex=LIGHT.lstrip("#"), line_hex=clr)
            rect(sl, x, 1.5, 3.0, 0.55, fill_hex=clr)
            txt(sl, phase, x+0.1, 1.55, 2.8, 0.35, size=14, bold=True, color=WHITE)
            txt(sl, period, x+0.1, 1.92, 2.8, 0.25, size=8.5, color=WHITE)
            for j, m in enumerate(milestones):
                rect(sl, x+0.15, 2.3+j*1.1, 2.7, 0.9, fill_hex=WHITE, line_hex=clr)
                txt(sl, m, x+0.25, 2.4+j*1.1, 2.5, 0.7, size=9, color=TEXT)

            # Arrow between phases
            if i < 3:
                txt(sl, "▶", x+3.0, 4.1, 0.25, 0.5, size=16, color=NAVY)

        # ── Slide 10: Expected Outcomes ────────────────────────────────────────
        sl = add_slide()
        title_bar(sl, "Expected Outcomes & ROI Model",
                  "Based on industry benchmarks and dataset analysis")

        outcomes = [
            ("Revenue Protected",  f"${kpis.revenue_at_risk * 0.28:,.0f}", "Base case scenario"),
            ("Churn Rate Target",  "< 15.0%",                              f"From {kpis.churn_rate_pct:.1f}% current"),
            ("Campaign ROI",       "2.1× – 3.5×",                          "$250K–$500K budget"),
            ("Satisfaction Target","≥ 7.5 / 10",                           f"From {kpis.avg_satisfaction_score:.1f} current"),
        ]
        for i, (label, val, note) in enumerate(outcomes):
            x = 0.4 + (i % 2) * 6.25
            y = 1.6 + (i // 2) * 2.5
            rect(sl, x, y, 5.8, 2.1, fill_hex=NAVY if i==0 else LIGHT.lstrip("#"),
                 line_hex=TEAL)
            txt(sl, label, x+0.2, y+0.15, 5.4, 0.4, size=10, color=WHITE if i==0 else MUTED)
            txt(sl, val,   x+0.2, y+0.6,  5.4, 0.9, size=26, bold=True,
                color=GOLD if i==0 else NAVY)
            txt(sl, note,  x+0.2, y+1.55, 5.4, 0.35, size=9, color=WHITE if i==0 else MUTED)

        txt(sl, "Disclaimer: Recovery projections are based on historical churn patterns and industry benchmarks. "
            "Actual results may vary based on execution quality and market conditions.",
            0.4, 6.9, 12.5, 0.4, size=7.5, color=MUTED, align=PP_ALIGN.CENTER)

        filename = f"mckinsey_retention_deck_{_ts()}.pptx"
        filepath = os.path.join(_ensure_reports_dir(), filename)
        prs.save(filepath)

        return ReportResponse(
            report_type="pptx",
            filename=filename,
            download_path=f"/reports/export/{filename}",
            generated_at=datetime.now(timezone.utc).isoformat(),
            size_bytes=os.path.getsize(filepath),
        )

    def _generate_markdown(self, req: ReportRequest) -> ReportResponse:
        kpis = self._ds.kpis()
        now  = datetime.now().strftime("%Y-%m-%d %H:%M UTC")
        md = f"""# Customer Retention Strategic Report\n_{now}_\n
## Executive Summary\nChurn: {kpis.churn_rate_pct:.1f}% | Revenue at Risk: ${kpis.revenue_at_risk:,.0f} | Customers: {kpis.total_customers:,}\n
## KPIs\n| Metric | Value |\n|--------|-------|\n| Churn Rate | {kpis.churn_rate_pct:.1f}% |\n| Revenue at Risk | ${kpis.revenue_at_risk:,.0f} |\n| Avg CLV | ${kpis.avg_clv:,.0f} |\n"""
        filename = f"report_{_ts()}.md"
        filepath = os.path.join(_ensure_reports_dir(), filename)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(md)
        return ReportResponse(
            report_type="markdown", filename=filename,
            download_path=f"/reports/export/{filename}",
            generated_at=datetime.now(timezone.utc).isoformat(),
            size_bytes=os.path.getsize(filepath),
        )

    def generate_markdown_report(self, req: ReportRequest) -> ReportResponse:
        return self._generate_markdown(req)

    @staticmethod

    @staticmethod
    def _risk_label(churn_pct: float) -> str:
        return _risk_label(churn_pct)
